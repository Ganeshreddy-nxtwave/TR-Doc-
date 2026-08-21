"""Read the corpus and resolve where a session sits in the course."""
import hashlib
import re
from pathlib import Path
from urllib.parse import urlparse

import yaml


def load_config(path="config.yaml"):
    return yaml.safe_load(Path(path).read_text(encoding="utf-8"))


def load_curriculum(path="curriculum.yaml"):
    data = yaml.safe_load(Path(path).read_text(encoding="utf-8"))
    sessions = data.get("sessions") or []
    if not sessions:
        raise SystemExit("curriculum.yaml has no sessions. Fill it in first.")
    # Sort by position, so an inserted session appended at the end of the file
    # still resolves to its real place instead of looking like the last one.
    ordered = sorted(sessions, key=sort_key)
    if ordered != sessions:
        print(f"Note: {path} was not in position order; using sorted order.")
    return data.get("course", "Unknown course"), ordered


def trusted_domains(path="sources.yaml"):
    return yaml.safe_load(Path(path).read_text(encoding="utf-8"))["trusted"]


def order_value(rec):
    """The ordering key. Tracker-derived curricula use `seq`; older ones use
    `unit_number`. Tracker order is preferred because LMS unit numbers collide
    and skip (they count quizzes and labs too)."""
    v = rec.get("seq")
    return str(v) if v not in (None, "") else str(rec.get("unit_number"))


def resolve_position(sessions, after):
    """Return (previous, next) for a session inserted after position `after`.

    `after` is matched as a string, so "12" and "12.5" both work and nothing
    downstream gets renumbered.
    """
    idx = next(
        (i for i, s in enumerate(sessions) if order_value(s) == str(after)),
        None,
    )
    if idx is None:
        have = ", ".join(order_value(s) for s in sessions)
        raise SystemExit(f"No session at position {after!r}. Have: {have}")
    nxt = sessions[idx + 1] if idx + 1 < len(sessions) else None
    return sessions[idx], nxt


def resolve_placement(sessions, position=None, after=None):
    """Where a NEW session goes. Returns (prev, nxt); either may be None.

    "first" -> nothing before it, sessions[0] after it
    "last"  -> sessions[-1] before it, nothing after
    "between" (or a bare `after`) -> delegates to resolve_position
    """
    if position == "first":
        return None, sessions[0]
    if position == "last":
        return sessions[-1], None
    if after in (None, ""):
        raise SystemExit(
            "Need either --position first/last or --after <position>."
        )
    return resolve_position(sessions, after)


def resolve_target(sessions, at):
    """The session being revamped or repurposed, plus its neighbours.

    Returns (target, prev, nxt). The target keeps its place, so `prev` is what
    comes before it -- not the target itself, which is what resolve_position
    returns for an insertion.
    """
    idx = next(
        (i for i, s in enumerate(sessions) if order_value(s) == str(at)),
        None,
    )
    if idx is None:
        have = ", ".join(order_value(s) for s in sessions)
        raise SystemExit(f"No session at position {at!r}. Have: {have}")
    prev = sessions[idx - 1] if idx > 0 else None
    nxt = sessions[idx + 1] if idx + 1 < len(sessions) else None
    return sessions[idx], prev, nxt


GSLIDES_RE = re.compile(r"docs\.google\.com/presentation/d/([a-zA-Z0-9_-]+)")
GDOCS_RE = re.compile(r"docs\.google\.com/document/d/([a-zA-Z0-9_-]+)")


def export_url(url):
    """Map a share link to a direct-download URL. Returns (url, suffix).

    Google Slides/Docs export endpoints need no auth, but ONLY if link-sharing
    is on. Anything else is passed through untouched.
    """
    m = GSLIDES_RE.search(url)
    if m:
        return (f"https://docs.google.com/presentation/d/{m.group(1)}/export/pptx",
                ".pptx")
    m = GDOCS_RE.search(url)
    if m:
        return (f"https://docs.google.com/document/d/{m.group(1)}/export?format=md",
                ".md")
    return url, (Path(urlparse(url).path).suffix or ".bin")


def fetch(url, cache_dir=".cache", suffix=None):
    """Download a source URL to a cache file and return its local path.

    A private Google file returns Google's HTML sign-in page instead of the file,
    so that case is detected and reported rather than ingested as garbage.
    """
    import urllib.error
    import urllib.request

    url, derived = export_url(url)
    suffix = suffix or derived

    cache = Path(cache_dir)
    cache.mkdir(parents=True, exist_ok=True)
    dest = cache / (hashlib.sha256(url.encode()).hexdigest()[:16] + suffix)
    if dest.exists():
        return dest

    req = urllib.request.Request(url, headers={"User-Agent": "tr-doc-generator"})
    ctype, data = "", None
    # Fetching ~90 decks in a run makes a transient reset near-certain, so retry
    # the connection-level failures. An HTTP status is a real answer: never retry.
    for attempt in range(1, 4):
        try:
            with urllib.request.urlopen(req, timeout=60) as r:
                ctype = (r.headers.get("Content-Type") or "").lower()
                data = r.read()
            break
        except urllib.error.HTTPError as e:
            raise SystemExit(
                f"Could not fetch {url} (HTTP {e.code}). If this is a Google file, "
                "it is probably not shared -- turn on link-sharing or download a "
                "copy."
            )
        except (urllib.error.URLError, ConnectionError, TimeoutError, OSError) as e:
            reason = getattr(e, "reason", e)
            if attempt == 3:
                raise SystemExit(
                    f"Could not reach {url} after 3 attempts: {reason}")
            print(f"  retry {attempt}/2 after {type(e).__name__}: {reason}")

    if suffix == ".pptx" and not data[:2] == b"PK":
        raise SystemExit(
            f"{url} did not return a .pptx file (got {ctype or 'unknown type'}). "
            "A Google Slides link that is not shared returns a sign-in page "
            "instead of the deck. Turn on link-sharing, or download a copy into "
            "corpus/ppts/."
        )

    dest.write_bytes(data)
    print(f"  fetched {url} -> {dest}")
    return dest


def resolve_source(path_or_url, cache_dir=".cache"):
    """Accept a local path or an http(s) URL; return a local Path."""
    if not path_or_url:
        return None
    s = str(path_or_url)
    if s.startswith(("http://", "https://")):
        return fetch(s, cache_dir)
    p = Path(s)
    return p if p.exists() else None


def read_text_file(path):
    """Read a TR doc. Markdown/text direct; docx and pdf via optional deps."""
    p = resolve_source(path)
    if p is None:
        return None
    suffix = p.suffix.lower()
    if suffix in (".md", ".txt", ".markdown", ""):
        return p.read_text(encoding="utf-8", errors="replace")
    if suffix == ".docx":
        try:
            import docx
        except ImportError:
            raise SystemExit(f"{p} is .docx -- run: pip install python-docx")
        return "\n".join(par.text for par in docx.Document(str(p)).paragraphs)
    if suffix == ".pdf":
        try:
            from pypdf import PdfReader
        except ImportError:
            raise SystemExit(f"{p} is .pdf -- run: pip install pypdf")
        return "\n".join(pg.extract_text() or "" for pg in PdfReader(str(p)).pages)
    raise SystemExit(f"Do not know how to read {p} ({suffix})")


def read_pptx(path):
    """Flatten a deck to text, one labelled block per slide. Accepts a URL."""
    from pptx import Presentation

    local = resolve_source(path)
    if local is None:
        return None
    out = []
    for n, slide in enumerate(Presentation(str(local)).slides, 1):
        lines = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        block = f"--- Slide {n} ---\n" + "\n".join(lines)
        if notes:
            block += f"\n[Speaker notes] {notes}"
        out.append(block)
    return "\n\n".join(out)


def read_links_file(path):
    """One URL per line, `#` for comments. Blank lines ignored."""
    p = Path(path)
    if not p.exists():
        return []
    out = []
    for line in p.read_text(encoding="utf-8").splitlines():
        line = line.split("#", 1)[0].strip()
        if line:
            out.append(line)
    return out


def collect_corpus(ppt_dir, doc_dir):
    """Everything on disk, for the one-time style distillation.

    Decks may be local .pptx files in ppt_dir, or URLs listed one per line in
    ppt_dir/links.txt -- those are downloaded and cached on first use.
    """
    items = []
    for p in sorted(Path(ppt_dir).glob("*.pptx")):
        if p.name.startswith("~$"):  # Office lock file
            continue
        items.append((p, read_pptx(p)))
    for url in read_links_file(Path(ppt_dir) / "links.txt"):
        text = read_pptx(url)
        if text:
            items.append((Path(url), text))
    for p in sorted(Path(doc_dir).rglob("*")) if Path(doc_dir).exists() else []:
        if p.is_file() and not p.name.startswith("~$"):
            text = read_text_file(p)
            if text:
                items.append((p, text))
    return items


# --- curriculum extraction -------------------------------------------------
# The docs carry a machine-readable header, so the curriculum is derived from
# them rather than typed by hand. Conflicts are reported, never resolved.

HEADER_RE = {
    "title": re.compile(r"^#\s+(.+?)\s*$", re.MULTILINE),
    "course": re.compile(r"\*\*Course:\*\*\s*(.+?)\s*(?:\||$)", re.MULTILINE),
    "topic": re.compile(r"\*\*Topic:\*\*\s*(.+?)\s*(?:\||$)", re.MULTILINE),
    "unit_id": re.compile(r"\*\*Unit ID:\*\*\s*`?([0-9a-fA-F]{8,})`?", re.MULTILINE),
    "unit_number": re.compile(r"\*\*Unit Number:\*\*\s*([0-9.]+)", re.MULTILINE),
}


def parse_header(text):
    """Pull the header block out of a TR doc. Missing fields come back None."""
    head = text[:2000]
    out = {}
    for key, rx in HEADER_RE.items():
        m = rx.search(head)
        out[key] = m.group(1).strip() if m else None
    return out


def scan_docs(doc_dir):
    """Every TR doc under doc_dir, with its parsed header and content hash."""
    records = []
    for p in sorted(Path(doc_dir).rglob("*.md")):
        text = read_text_file(p)
        if not text:
            continue
        rec = parse_header(text)
        rec["path"] = p.as_posix()
        rec["sha"] = hashlib.sha256(text.encode("utf-8")).hexdigest()[:12]
        rec["chars"] = len(text)
        records.append(rec)
    return records


def curriculum_issues(records):
    """Everything a human must decide. Reported, never auto-resolved."""
    issues = []

    by_sha = {}
    for r in records:
        by_sha.setdefault(r["sha"], []).append(r["path"])
    for sha, paths in by_sha.items():
        if len(paths) > 1:
            issues.append("DUPLICATE CONTENT (same file saved under different "
                          "names -- one of these is the wrong doc):\n    "
                          + "\n    ".join(paths))

    by_unit = {}
    for r in records:
        if r["course"] and r["unit_number"]:
            by_unit.setdefault((r["course"], r["unit_number"]), []).append(r)
    for (course, num), rs in sorted(by_unit.items()):
        if len({r["unit_id"] for r in rs}) > 1:
            issues.append(f"UNIT NUMBER CONFLICT: {course} unit {num} claimed by "
                          f"{len(rs)} different Unit IDs:\n    "
                          + "\n    ".join(f"{r['unit_id']}  {r['path']}" for r in rs))

    for r in records:
        missing = [k for k in ("course", "topic", "unit_id", "unit_number")
                   if not r[k]]
        if missing:
            issues.append(f"MISSING HEADER FIELDS {missing} in {r['path']}")

    return issues


def sort_key(rec):
    """Order by seq, or unit_number when there is no seq. Works for scanned docs
    and for curriculum entries alike."""
    where = rec.get("path") or rec.get("tr_doc") or ""
    try:
        return (0, float(order_value(rec) or 0), where)
    except (TypeError, ValueError):
        return (1, 0.0, where)


def render_curriculum(records, course_filter=None):
    """Emit curriculum.yaml text for one course, in teaching order."""
    rows = [r for r in records
            if r["course"] and (not course_filter or r["course"] == course_filter)]
    if not rows:
        raise SystemExit(f"No docs found for course {course_filter!r}")
    course = course_filter or rows[0]["course"]
    rows.sort(key=sort_key)

    lines = [
        "# Generated by `python -m tr curriculum`. Regenerate after adding docs.",
        "# Sessions are in teaching order. Insert a new session by adding an entry",
        "# with a decimal unit_number (e.g. \"33.5\") -- nothing downstream renumbers.",
        "",
        f"course: {course!r}",
        "",
        "sessions:",
    ]
    for r in rows:
        lines += [
            f"  - unit_id: {(r['unit_id'] or 'UNKNOWN')!r}",
            f"    unit_number: {(r['unit_number'] or '0')!r}",
            f"    title: {(r['title'] or 'UNTITLED')!r}",
            f"    topic: {(r['topic'] or '')!r}",
            f"    tr_doc: {r['path']}",
        ]
    return "\n".join(lines) + "\n"


# --- status-tracker import -------------------------------------------------
# The tracker CSV is authoritative for teaching ORDER and for PPT links. The doc
# headers are authoritative for Unit ID and LMS unit number. This merges them and
# reports whatever it could not match, rather than guessing.

import csv
import difflib

TRACKER_COLS = ("Course", "Topic Name", "Session Name", "PPT Link")


def read_tracker(path):
    """Rows from the status tracker, with merged cells forward-filled.

    The sheet has two header rows; the real one is the second. Course and Topic
    are merged cells in the sheet, so they arrive blank on continuation rows.
    """
    with Path(path).open(encoding="utf-8-sig", newline="") as fh:
        rows = list(csv.reader(fh))
    header_idx = next(
        (i for i, r in enumerate(rows[:5])
         if all(c in [x.strip() for x in r] for c in TRACKER_COLS)),
        None,
    )
    if header_idx is None:
        raise SystemExit(
            f"{path} does not look like the status tracker. Expected a header row "
            f"containing {', '.join(TRACKER_COLS)}."
        )
    hdr = [h.strip() for h in rows[header_idx]]

    course = topic = ""
    out = []
    for r in rows[header_idx + 1:]:
        if not any(c.strip() for c in r):
            continue
        d = dict(zip(hdr, r))
        course = d.get("Course", "").strip() or course
        topic = d.get("Topic Name", "").strip() or topic
        session = d.get("Session Name", "").strip()
        if not session:
            continue
        # Sheets' CSV export keeps only the DISPLAY TEXT of a hyperlinked cell,
        # so a "PPT Link" that is not a URL is a lost link, not a usable one.
        # A "PPT URL" column, if present, is a recovered-link column and wins.
        cell = (d.get("PPT URL") or "").strip() or d.get("PPT Link", "").strip()
        is_url = cell.startswith(("http://", "https://"))
        out.append({
            "row": len(out) + 1,
            "course": course,
            "topic": topic,
            "title": session,
            "ppt": cell if is_url else None,
            "ppt_label": None if is_url else (cell or None),
        })
    return out


def norm_title(s):
    """Collapse a title to comparable words. '| Part 1' and '- 1' both survive."""
    s = re.sub(r"[|\-_/]+", " ", (s or "").lower())
    s = re.sub(r"[^a-z0-9 ]+", " ", s)
    return " ".join(s.split())


STOPWORDS = {"a", "an", "the", "with", "in", "on", "for", "to", "of", "and",
             "using", "your", "part", "own", "build", "1", "2", "3"}


def content_words(s):
    return frozenset(w for w in s.split() if w not in STOPWORDS)


def title_score(a, b):
    """Similarity of two normalised titles. Three signals, best one wins.

    1. Character edit ratio -- the baseline.
    2. Containment: 'integrating mcp' inside 'integrating mcp servers in
       langchain agents' scores low on edit ratio but is a strong match, so it
       is floored at 0.9. Both sides must be substantial so a short generic
       title cannot swallow an unrelated one.
    3. Same content words in any order: 'Building a Memory Agent' vs 'Building
       an Agent with Memory'. Identical word sets floor at 0.95; otherwise the
       word-set overlap is used, which beats edit distance on reordered titles.
    """
    ratio = difflib.SequenceMatcher(None, a, b).ratio()

    short, long = sorted((a, b), key=len)
    if len(short) >= 12 and short in long:
        return max(ratio, 0.9)

    wa, wb = content_words(a), content_words(b)
    if wa and wb:
        if wa == wb:
            return max(ratio, 0.95)
        ratio = max(ratio, len(wa & wb) / len(wa | wb))
    return ratio


def match_docs(tracker_rows, doc_records, cutoff=0.8, course=None):
    """Attach tr_doc, unit_id and unit_number to tracker rows by title match.

    `course` restricts the candidate pool to docs whose header names that course,
    so a generic title like 'Your Learning Journey' cannot match the wrong
    course's doc. If no doc carries that course name, the pool is left unfiltered.

    Returns (rows, unmatched_docs). Each row gains 'tr_doc', 'unit_id',
    'unit_number' and 'match_score'. A doc is used at most once.
    """
    if course:
        scoped = [r for r in doc_records
                  if (r.get("course") or "").strip().lower() == course.strip().lower()]
        if not scoped:
            print(f"  No TR doc header names the course {course!r}. Matching "
                  "nothing rather than borrowing another course's docs -- pass "
                  "--doc-course if the header spells it differently.")
        doc_records = scoped

    pool = {r["path"]: r for r in doc_records}
    keys = {p: norm_title(r.get("title")) for p, r in pool.items()}

    for row in tracker_rows:
        want = norm_title(row["title"])
        best, best_score = None, 0.0
        for path, key in keys.items():
            if path not in pool:
                continue
            score = title_score(want, key)
            if score > best_score:
                best, best_score = path, score
        if best and best_score >= cutoff:
            rec = pool.pop(best)
            row.update(tr_doc=rec["path"], unit_id=rec.get("unit_id"),
                       unit_number=rec.get("unit_number"),
                       match_score=round(best_score, 3))
        else:
            row.update(tr_doc=None, unit_id=None, unit_number=None,
                       match_score=round(best_score, 3))
    return tracker_rows, list(pool.values())


def render_tracker_curriculum(rows, course_label):
    """Emit curriculum.yaml where `seq` is the ordering key from the tracker."""
    lines = [
        "# Generated by `python -m tr curriculum --tracker <csv>`.",
        "# `seq` is the ordering key and comes from tracker row order.",
        "# Insert a session with a decimal seq (e.g. 12.5) -- nothing renumbers.",
        "# `unit_number` and `unit_id` are the LMS values, present only where a TR",
        "# doc was matched. Null means the tool must ask rather than guess.",
        "",
        f"course: {course_label!r}",
        "",
        "sessions:",
    ]
    for i, r in enumerate(rows, 1):
        lines += [
            f"  - seq: '{i}'",
            f"    title: {r['title']!r}",
            f"    topic: {(r.get('topic') or '')!r}",
        ]
        lines.append(f"    unit_id: {r['unit_id']!r}" if r.get("unit_id")
                     else "    unit_id: null")
        lines.append(f"    unit_number: {r['unit_number']!r}" if r.get("unit_number")
                     else "    unit_number: null")
        lines.append(f"    tr_doc: {r['tr_doc']}" if r.get("tr_doc")
                     else "    tr_doc: null")
        if r.get("ppt"):
            lines.append(f"    ppt: {r['ppt']}")
        else:
            lines.append("    ppt: null")
            if r.get("ppt_label"):
                lines.append(f"    # deck exists but the CSV export kept only its "
                             f"label: {r['ppt_label']!r}")
    return "\n".join(lines) + "\n"
