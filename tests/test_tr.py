"""Offline self-check. Run: python tests/test_tr.py

Covers the logic that can be wrong silently: snippet classification, real
execution and splicing, neighbour resolution, and domain trust matching.
No network, no API key, no LLM calls.
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from tr import corpus, generate, research, runner
from tr.cli import slugify


def test_classification():
    assert runner.is_live("from openai import OpenAI\nc = OpenAI()")
    assert runner.is_live("import anthropic")
    assert runner.is_live("r = requests.get(url)")
    assert runner.is_live("name = input('who? ')")
    assert runner.is_live("m = AutoModel.from_pretrained('bert')")
    assert not runner.is_live("from pydantic import BaseModel\nprint(1)")
    assert not runner.is_live("import json\nprint(json.dumps({'a': 1}))")


def test_runs_offline_snippet_and_inserts_real_output():
    md = "Some prose.\n\n```python\nprint(2 + 2)\n```\n\nMore prose.\n"
    out, results = runner.verify(md, timeout=20)
    assert results[0]["status"] == "ok", results
    assert "Output of a real run" in out
    assert "```text\n4\n```" in out, out
    assert "More prose." in out


def test_live_snippet_is_marked_not_run():
    md = "```python\nfrom openai import OpenAI\nprint(OpenAI())\n```\n"
    out, results = runner.verify(md, timeout=20)
    assert results[0]["status"] == "live"
    assert "[UNVERIFIED]" in out
    assert "paste the output of your run" in out


def test_failing_snippet_shows_the_real_error():
    md = "```python\nraise ValueError('boom')\n```\n"
    out, results = runner.verify(md, timeout=20)
    assert results[0]["status"] == "error"
    assert "boom" in out
    assert "[UNVERIFIED]" in out


def test_missing_dependency_is_not_reported_as_success():
    md = "```python\nimport definitely_not_a_real_module_xyz\n```\n"
    out, results = runner.verify(md, timeout=20)
    assert results[0]["status"] == "missing_dep", results
    assert "definitely_not_a_real_module_xyz" in out


def test_multiple_snippets_keep_their_order():
    md = "```python\nprint('first')\n```\ntext\n```python\nprint('second')\n```\n"
    out, results = runner.verify(md, timeout=20)
    assert len(results) == 2
    assert out.index("first") < out.index("text") < out.index("second")


def test_resolve_position():
    sessions = [
        {"unit_number": "3", "title": "Prompting", "tr_doc": "a.md"},
        {"unit_number": "4", "title": "Function Calling", "tr_doc": "b.md"},
        {"unit_number": "5", "title": "Agents", "tr_doc": "c.md"},
    ]
    prev, nxt = corpus.resolve_position(sessions, "4")
    assert prev["title"] == "Function Calling"
    assert nxt["title"] == "Agents"

    prev, nxt = corpus.resolve_position(sessions, "5")
    assert nxt is None, "last session must yield no next"

    try:
        corpus.resolve_position(sessions, "99")
    except SystemExit as e:
        assert "99" in str(e)
    else:
        raise AssertionError("unknown unit_number must fail loudly")


def test_inserted_session_resolves_even_when_appended_out_of_order():
    """A session added at the end of the file must still sort into position."""
    import tempfile
    yml = """course: 'C'
sessions:
  - unit_number: '43'
    title: 'MCP'
    tr_doc: a.md
  - unit_number: '45'
    title: 'CrewAI'
    tr_doc: b.md
  - unit_number: '43.5'
    title: 'Inserted'
    tr_doc: c.md
"""
    with tempfile.TemporaryDirectory() as d:
        p = Path(d) / "cur.yaml"
        p.write_text(yml, encoding="utf-8")
        _, sessions = corpus.load_curriculum(p)
    assert [s["unit_number"] for s in sessions] == ["43", "43.5", "45"]

    prev, nxt = corpus.resolve_position(sessions, "43.5")
    assert prev["title"] == "Inserted"
    assert nxt["title"] == "CrewAI", "inserted session must not look like the last one"

    prev, nxt = corpus.resolve_position(sessions, "43")
    assert nxt["title"] == "Inserted", "unit 43's next is now the inserted session"


def test_parse_header():
    doc = (
        "# Introduction to AI Agents\n\n"
        "**Course:** Generative AI  \n"
        "**Topic:** Building AI Agents  \n"
        "**Unit ID:** `2e566865957343ed9579b659d494dc63` | **Unit Number:** 33\n"
    )
    h = corpus.parse_header(doc)
    assert h["title"] == "Introduction to AI Agents"
    assert h["course"] == "Generative AI"
    assert h["topic"] == "Building AI Agents"
    assert h["unit_id"] == "2e566865957343ed9579b659d494dc63"
    assert h["unit_number"] == "33"

    empty = corpus.parse_header("just some text, no header")
    assert all(v is None for v in empty.values())


def test_curriculum_issues_detects_both_problems():
    records = [
        {"course": "C", "topic": "t", "unit_id": "AAA", "unit_number": "33",
         "path": "x.md", "sha": "same", "chars": 10},
        {"course": "C", "topic": "t", "unit_id": "BBB", "unit_number": "33",
         "path": "y.md", "sha": "same", "chars": 10},
    ]
    issues = corpus.curriculum_issues(records)
    joined = "\n".join(issues)
    assert "DUPLICATE CONTENT" in joined
    assert "UNIT NUMBER CONFLICT" in joined

    clean = corpus.curriculum_issues([
        {"course": "C", "topic": "t", "unit_id": "AAA", "unit_number": "1",
         "path": "x.md", "sha": "a", "chars": 10}])
    assert clean == []


def test_google_share_links_become_export_urls():
    url, sfx = corpus.export_url(
        "https://docs.google.com/presentation/d/1AbC-dEf_9/edit#slide=id.p1")
    assert url == "https://docs.google.com/presentation/d/1AbC-dEf_9/export/pptx"
    assert sfx == ".pptx"

    url, sfx = corpus.export_url("https://docs.google.com/document/d/XyZ123/edit")
    assert url.endswith("/export?format=md")
    assert sfx == ".md"

    # a direct file URL is passed through untouched
    url, sfx = corpus.export_url("https://cdn.example.com/decks/session-04.pptx")
    assert url == "https://cdn.example.com/decks/session-04.pptx"
    assert sfx == ".pptx"

    # unknown extension still gets a usable suffix
    _, sfx = corpus.export_url("https://example.com/download?id=7")
    assert sfx == ".bin"


def test_links_file_parsing():
    import tempfile
    with tempfile.TemporaryDirectory() as d:
        p = Path(d) / "links.txt"
        p.write_text(
            "# a comment\n"
            "\n"
            "https://example.com/a.pptx\n"
            "  https://example.com/b.pptx   # trailing comment\n"
            "# https://example.com/disabled.pptx\n",
            encoding="utf-8")
        assert corpus.read_links_file(p) == [
            "https://example.com/a.pptx", "https://example.com/b.pptx"]
    assert corpus.read_links_file(Path(d) / "gone.txt") == []


def test_resolve_source_returns_none_for_missing_local_file():
    assert corpus.resolve_source(None) is None
    assert corpus.resolve_source("") is None
    assert corpus.resolve_source("definitely/not/here.md") is None
    assert corpus.resolve_source(__file__) is not None


def test_read_pptx_extracts_slides_and_speaker_notes():
    try:
        from pptx import Presentation
    except ImportError:
        print("      (skipped: python-pptx not installed)")
        return
    import tempfile

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "What breaks without tools"
    slide.placeholders[1].text = "The model cannot read today's weather"
    slide.notes_slide.notes_text_frame.text = "Ask the class before revealing"

    with tempfile.TemporaryDirectory() as d:
        p = Path(d) / "deck.pptx"
        prs.save(str(p))
        text = corpus.read_pptx(p)

    assert "--- Slide 1 ---" in text
    assert "What breaks without tools" in text
    assert "The model cannot read today's weather" in text
    assert "[Speaker notes] Ask the class before revealing" in text, \
        "speaker notes carry the teaching intent and must not be dropped"

    assert corpus.read_pptx("no/such/deck.pptx") is None


TRACKER_CSV = (
    " ,,,,,,VIDEO TRACKER\n"
    "Course,Topic Name,Session Name,PPT Link,Owner\n"
    ",,Your Learning Journey,https://docs.google.com/presentation/d/AAA/edit,me\n"
    ",,Gen AI in Action,,me\n"
    "Building LLM Applications,Overview,Introduction to Langchain,Intro to LangChain,me\n"
    ",,Building Memory Agent,https://example.com/d.pptx,me\n"
    ",,,,\n"
)


def test_read_tracker_forward_fills_and_separates_lost_links():
    import tempfile
    with tempfile.TemporaryDirectory() as d:
        p = Path(d) / "t.csv"
        p.write_text(TRACKER_CSV, encoding="utf-8")
        rows = corpus.read_tracker(p)

    assert len(rows) == 4, "blank rows and the sub-header must not become sessions"
    assert rows[0]["ppt"] == "https://docs.google.com/presentation/d/AAA/edit"
    assert rows[0]["ppt_label"] is None

    # empty cell: neither a url nor a lost link
    assert rows[1]["ppt"] is None and rows[1]["ppt_label"] is None

    # display text only -- the URL was dropped by the CSV export, so it must NOT
    # land in `ppt`, or the fetcher would try to download a sentence
    assert rows[2]["ppt"] is None
    assert rows[2]["ppt_label"] == "Intro to LangChain"

    # course and topic forward-fill onto continuation rows
    assert rows[3]["course"] == "Building LLM Applications"
    assert rows[3]["topic"] == "Overview"
    assert rows[3]["ppt"] == "https://example.com/d.pptx"


def test_read_tracker_rejects_a_non_tracker_csv():
    import tempfile
    with tempfile.TemporaryDirectory() as d:
        p = Path(d) / "other.csv"
        p.write_text("a,b,c\n1,2,3\n", encoding="utf-8")
        try:
            corpus.read_tracker(p)
        except SystemExit as e:
            assert "status tracker" in str(e)
        else:
            raise AssertionError("a random CSV must be rejected, not half-parsed")


def test_title_score_handles_reorder_containment_and_near_misses():
    ns = corpus.norm_title
    # same content words, different order -- must match
    assert corpus.title_score(ns("Building a Memory Agent"),
                              ns("Building an Agent with Memory")) >= 0.9
    # one title contained in the other
    assert corpus.title_score(ns("Integrating MCP"),
                              ns("Integrating MCP Servers in LangChain Agents")) >= 0.9
    assert corpus.title_score(ns("Your Learning Journey"),
                              ns("Your Learning Journey | Building LLM Apps")) >= 0.9
    # different sessions that merely share an opening word -- must stay below cutoff
    assert corpus.title_score(ns("Introduction to Generative AI"),
                              ns("Introduction to AI Agents")) < 0.8
    assert corpus.title_score(ns("Running Models Locally"),
                              ns("Finetuning the LLMs")) < 0.8


def test_match_docs_refuses_to_borrow_another_courses_doc():
    docs = [{"path": "genai/02.md", "title": "Your Learning Journey",
             "course": "Generative AI", "unit_id": "A", "unit_number": "2"}]
    rows = [{"row": 1, "course": "AI for Finance", "topic": "",
             "title": "Your Learning Journey (AI for Finance)", "ppt": None}]

    matched, unmatched = corpus.match_docs(list(rows), docs, course="AI for Finance")
    assert matched[0]["tr_doc"] is None, \
        "no AI for Finance doc exists, so it must match nothing"
    assert len(unmatched) == 0

    # with the correct course it does match
    rows2 = [{"row": 1, "course": "", "topic": "", "title": "Your Learning Journey",
              "ppt": None}]
    matched, _ = corpus.match_docs(rows2, docs, course="Generative AI")
    assert matched[0]["tr_doc"] == "genai/02.md"
    assert matched[0]["unit_id"] == "A"


def test_match_docs_uses_each_doc_at_most_once():
    docs = [
        {"path": "a.md", "title": "Introduction to RAG Part 1", "course": "C",
         "unit_id": "1", "unit_number": "1"},
        {"path": "b.md", "title": "Introduction to RAG Part 2", "course": "C",
         "unit_id": "2", "unit_number": "2"},
    ]
    rows = [
        {"row": 1, "course": "C", "topic": "", "title": "Introduction to RAG Part 1",
         "ppt": None},
        {"row": 2, "course": "C", "topic": "", "title": "Introduction to RAG Part 2",
         "ppt": None},
    ]
    matched, unmatched = corpus.match_docs(rows, docs, course="C")
    assert {r["tr_doc"] for r in matched} == {"a.md", "b.md"}
    assert unmatched == []


def test_seq_is_the_ordering_key_when_present():
    sessions = [
        {"seq": "1", "title": "A"},
        {"seq": "2", "title": "B"},
        {"seq": "1.5", "title": "Inserted"},
    ]
    ordered = sorted(sessions, key=corpus.sort_key)
    assert [s["title"] for s in ordered] == ["A", "Inserted", "B"]

    prev, nxt = corpus.resolve_position(ordered, "1.5")
    assert prev["title"] == "Inserted" and nxt["title"] == "B"

    # falls back to unit_number when there is no seq
    legacy = [{"unit_number": "5", "title": "X"}, {"unit_number": "7", "title": "Y"}]
    prev, nxt = corpus.resolve_position(legacy, "5")
    assert nxt["title"] == "Y"


def test_order_value_and_describe_never_print_none():
    """Tracker curricula have a seq but often a null unit_number. Neither the
    ordering key nor the human label may leak 'None' into output."""
    from tr.cli import describe

    tracker_style = {"seq": "7", "unit_number": None,
                     "title": "Building a Trading Agent | Part 1"}
    assert corpus.order_value(tracker_style) == "7"
    label = describe(tracker_style)
    assert "None" not in label
    assert "7." in label and "no LMS unit" in label

    matched = {"seq": "13", "unit_number": "27", "title": "Mastering Image Generation"}
    assert corpus.order_value(matched) == "13", "seq wins over unit_number"
    assert "LMS unit 27" in describe(matched)

    legacy = {"unit_number": "27", "title": "X"}
    assert corpus.order_value(legacy) == "27"
    assert "None" not in describe(legacy)


class _FakeResponse:
    """Minimal stand-in for the urlopen context manager."""

    def __init__(self, data, ctype="application/octet-stream"):
        self._data = data
        self.headers = {"Content-Type": ctype}

    def __enter__(self):
        return self

    def __exit__(self, *a):
        return False

    def read(self):
        return self._data


def _with_fake_urlopen(fn):
    """Swap urllib.request.urlopen for the duration of a call."""
    import urllib.request
    real = urllib.request.urlopen
    urllib.request.urlopen = fn
    return real


def test_fetch_retries_transient_errors_but_not_http_status():
    import tempfile
    import urllib.error
    import urllib.request

    calls = {"n": 0}

    def flaky(req, timeout=None):
        calls["n"] += 1
        if calls["n"] < 3:
            raise ConnectionResetError("forcibly closed by the remote host")
        return _FakeResponse(b"PK\x03\x04payload")

    real = _with_fake_urlopen(flaky)
    try:
        with tempfile.TemporaryDirectory() as d:
            got = corpus.fetch("https://example.com/deck.pptx", cache_dir=d)
            assert got.read_bytes().startswith(b"PK")
        assert calls["n"] == 3, "must retry a reset, not give up on the first one"

        # an HTTP status is a real answer -- retrying it is pointless
        http_calls = {"n": 0}

        def forbidden(req, timeout=None):
            http_calls["n"] += 1
            raise urllib.error.HTTPError(req.full_url, 403, "Forbidden", {}, None)

        urllib.request.urlopen = forbidden
        with tempfile.TemporaryDirectory() as d:
            try:
                corpus.fetch("https://example.com/deck.pptx", cache_dir=d)
            except SystemExit as e:
                assert "403" in str(e)
            else:
                raise AssertionError("a 403 must fail, not silently succeed")
        assert http_calls["n"] == 1, "an HTTP status must not be retried"

        # persistent connection failure gives up with a clear message
        def always_dead(req, timeout=None):
            raise ConnectionResetError("still dead")

        urllib.request.urlopen = always_dead
        with tempfile.TemporaryDirectory() as d:
            try:
                corpus.fetch("https://example.com/deck.pptx", cache_dir=d)
            except SystemExit as e:
                assert "3 attempts" in str(e)
            else:
                raise AssertionError("must give up after the retry budget")
    finally:
        urllib.request.urlopen = real


def test_fetch_caches_and_does_not_refetch():
    import tempfile
    import urllib.request

    calls = {"n": 0}

    def once(req, timeout=None):
        calls["n"] += 1
        return _FakeResponse(b"PK\x03\x04payload")

    real = _with_fake_urlopen(once)
    try:
        with tempfile.TemporaryDirectory() as d:
            a = corpus.fetch("https://example.com/deck.pptx", cache_dir=d)
            b = corpus.fetch("https://example.com/deck.pptx", cache_dir=d)
        assert a == b
        assert calls["n"] == 1, "a cached URL must not be downloaded twice"
    finally:
        urllib.request.urlopen = real


SESSIONS = [
    {"seq": "1", "title": "Intro", "tr_doc": "a.md"},
    {"seq": "2", "title": "Middle", "tr_doc": "b.md"},
    {"seq": "3", "title": "Finale", "tr_doc": "c.md"},
]


def test_resolve_placement_first_last_between():
    prev, nxt = corpus.resolve_placement(SESSIONS, position="first")
    assert prev is None, "a first session has nothing before it"
    assert nxt["title"] == "Intro"

    prev, nxt = corpus.resolve_placement(SESSIONS, position="last")
    assert prev["title"] == "Finale"
    assert nxt is None, "a last session has nothing after it"

    prev, nxt = corpus.resolve_placement(SESSIONS, after="2")
    assert prev["title"] == "Middle" and nxt["title"] == "Finale"

    # neither position nor after is a usage error, not a silent default
    try:
        corpus.resolve_placement(SESSIONS)
    except SystemExit as e:
        assert "--position" in str(e) and "--after" in str(e)
    else:
        raise AssertionError("must refuse to guess the placement")


def test_resolve_placement_single_session_course():
    one = [{"seq": "1", "title": "Only", "tr_doc": "a.md"}]
    prev, nxt = corpus.resolve_placement(one, position="first")
    assert prev is None and nxt["title"] == "Only"
    prev, nxt = corpus.resolve_placement(one, position="last")
    assert prev["title"] == "Only" and nxt is None


def test_resolve_target_returns_neighbours_not_itself():
    target, prev, nxt = corpus.resolve_target(SESSIONS, "2")
    assert target["title"] == "Middle"
    assert prev["title"] == "Intro", \
        "prev must be the session BEFORE the target, unlike an insertion"
    assert nxt["title"] == "Finale"

    target, prev, nxt = corpus.resolve_target(SESSIONS, "1")
    assert prev is None and nxt["title"] == "Middle"

    target, prev, nxt = corpus.resolve_target(SESSIONS, "3")
    assert prev["title"] == "Middle" and nxt is None

    try:
        corpus.resolve_target(SESSIONS, "99")
    except SystemExit as e:
        assert "99" in str(e)
    else:
        raise AssertionError("an unknown position must fail loudly")


def test_mode_and_position_briefs_load_and_reject_unknowns():
    for name in ("new", "revamp", "repurpose"):
        assert len(generate.block("modes", name)) > 50
    for name in ("first", "middle", "last"):
        assert len(generate.block("positions", name)) > 50

    assert "CHANGES MADE" in generate.block("modes", "revamp")
    assert "REMOVE" in generate.block("modes", "repurpose").upper()
    assert "HOOK FOUNDATION" in generate.block("positions", "first")
    assert "no next session" in generate.block("positions", "last").lower()

    try:
        generate.block("modes", "nonsense")
    except SystemExit as e:
        assert "revamp" in str(e), "the error should list the valid names"
    else:
        raise AssertionError("an unknown brief must fail, not silently vanish")


def test_system_prompt_assembles_with_no_dangling_placeholders():
    """The real failure mode: a placeholder that is never filled ships to the
    model as literal '{{mode_brief}}' and the run silently loses its rules."""
    import re as _re

    filled = generate.fill(generate.prompt("tr_doc.md"), {
        "course": "Generative AI",
        "topic": "Structured Outputs",
        "previous_session": "12. Productivity Power-Up",
        "next_session": "13. Image Generation",
        "learner_profile": "knows Python basics",
        "session_produces": "working code",
        "mode_brief": generate.block("modes", "revamp"),
        "position_brief": generate.block("positions", "last"),
    })

    left = _re.findall(r"\{\{(\w+)\}\}", filled)
    assert not left, f"unfilled placeholders reached the model: {left}"

    assert "MODE: REVAMP" in filled
    assert "POSITION: this is the LAST session" in filled
    assert "CHANGES MADE" in filled
    assert "THE SEVEN RULES" in filled, "the base rules must survive injection"
    # the run-specific blocks must be able to override the general rules
    assert "they win" in filled

    # a first-position run must not carry middle-position instructions
    first = generate.fill(generate.prompt("tr_doc.md"), {
        "mode_brief": generate.block("modes", "new"),
        "position_brief": generate.block("positions", "first"),
    })
    assert "HOOK FOUNDATION" in first
    assert "POSITION: this session sits BETWEEN" not in first


class _StubClient:
    """Records the kwargs of each completions call instead of making one."""

    def __init__(self):
        self.calls = []
        outer = self

        class _Completions:
            def create(self, **kwargs):
                outer.calls.append(kwargs)

                class _M:
                    content = "# Doc\nbody"
                    annotations = None

                class _C:
                    message = _M()

                class _R:
                    choices = [_C()]

                return _R()

        class _Chat:
            completions = _Completions()

        self.chat = _Chat()


def test_writer_call_sends_an_explicit_token_budget():
    """Without max_tokens the provider default applies, which caps the doc far
    below the ~1,000 lines a real session needs. This was the single biggest
    cause of thin output."""
    cl = _StubClient()
    ctx = {
        "course": "C", "topic": "T", "prev_title": "1. Prev", "next_title": "3. Next",
        "prev_doc": "prev", "next_doc": "next", "research": "notes",
        "style": "style guide", "learner_profile": "x", "session_produces": "y",
        "max_output_tokens": 32000,
    }
    generate.write_doc(cl, "some/model", ctx, "answers")

    assert len(cl.calls) == 1
    assert cl.calls[0].get("max_tokens") == 32000, \
        "the writer must cap output explicitly, not inherit the provider default"

    # and config.yaml actually carries a value in that range
    cfg = corpus.load_config()
    assert cfg.get("max_output_tokens", 0) >= 16000, \
        "config max_output_tokens is too low for a full doc"

    # falls back to a usable budget rather than None if config omits it
    cl2 = _StubClient()
    generate.write_doc(cl2, "some/model", {**ctx, "max_output_tokens": None},
                       "answers")
    assert cl2.calls[0].get("max_tokens", 0) >= 16000


def test_depth_rules_are_calibrated_to_the_reference_doc():
    """The prompt is calibrated to the AI-agents reference doc: ~12 sections with
    deep insides. Guards against drifting thin (the first real output) or heavy
    (the 1,685-line doc's meta-scaffolding)."""
    t = " ".join(generate.prompt("tr_doc.md").split())

    # the nine patterns adopted from the reference
    for required in (
        "named beats",                    # hook structure
        "escalate",                       # the harder, unpatchable case
        "what the learner will actually write today",   # components table
        "piece by piece",                 # incremental assembly
        "what breaks without it",         # the reason per piece
        "properties after building it",   # post-build behaviour
        "named scenarios",                # several runs
        "negative case",                  # the boundary case
        "inline, where the choice arises",  # design decisions
        "discriminator",                  # Try It Yourself close
        "900 to 1,200 lines",             # explicit calibration
    ):
        assert required in t, f"depth rule missing: {required!r}"

    # the meta-scaffolding that made the other reference too heavy
    for banned in ("verify-flags table", "an agenda", "a coverage ledger",
                   "a checkpoint", "a key-takeaways list"):
        assert banned in t, f"the ban on {banned!r} should be stated explicitly"
    assert "Do NOT add:" in t

    # What's Next survives, since the reference doc wrongly omits it
    assert "**What's Next** -- required, never omitted" in t

    # house style the reference doc drops must stay required
    assert "<MultiLineWarning>" in t
    assert 'target="_blank"' in t

    # every trailing report block the code splits off must be documented here
    for blk in ("SOURCE ISSUES", "OPEN MARKERS", "CHANGES MADE",
                "ADDED BEYOND SCOPE"):
        assert blk in t, f"{blk} not documented in OUTPUT FORMAT"
        assert generate.TRAILING_RE.search(f"\n{blk}\n"), \
            f"{blk} documented but the splitter does not match it"


def test_job_line_describes_the_actual_run():
    assert "FIRST session" in generate.job_line(
        {"mode": "new", "position": "first"})
    assert "LAST session" in generate.job_line(
        {"mode": "new", "position": "last"})
    assert "after position 12" in generate.job_line(
        {"mode": "new", "position": "middle", "after": "12"})
    line = generate.job_line(
        {"mode": "revamp", "at": "12", "target_title": "12. Prompting"})
    assert "REVAMP" in line and "12. Prompting" in line


def test_execute_off_runs_nothing_and_marks_everything_unverified():
    """The hosted path. Nothing may execute, and no output may look verified."""
    from tr import pipeline  # noqa: F401  (import-ability is part of the check)

    md = ("```python\nprint(2 + 2)\n```\n\ntext\n\n"
          "```python\nfrom openai import OpenAI\n```\n")
    out, results = runner.verify(md, timeout=20, execute=False)

    assert [r["status"] for r in results] == ["not_executed", "not_executed"]
    assert "Output of a real run" not in out, \
        "nothing ran, so nothing may claim to be a real run"
    assert out.count("[UNVERIFIED]") == 2
    assert "```text\n4\n```" not in out, "the offline snippet must NOT have run"
    assert "text" in out, "surrounding prose must survive"

    summary = runner.summarise(results)
    assert "execution off here" in summary

    # and the default is still to execute
    out_on, results_on = runner.verify("```python\nprint(2 + 2)\n```\n", timeout=20)
    assert results_on[0]["status"] == "ok"
    assert "```text\n4\n```" in out_on


def test_pipeline_resolve_matches_the_cli_shapes():
    from tr import pipeline

    target, prev, nxt = pipeline.resolve(SESSIONS, "new", position="first")
    assert target is None and prev is None and nxt["title"] == "Intro"

    target, prev, nxt = pipeline.resolve(SESSIONS, "new", position="last")
    assert target is None and nxt is None and prev["title"] == "Finale"

    # "middle" and "between" both mean the same thing to callers
    for word in ("middle", "between"):
        target, prev, nxt = pipeline.resolve(SESSIONS, "new", position=word,
                                             after="2")
        assert target is None
        assert prev["title"] == "Middle" and nxt["title"] == "Finale"

    target, prev, nxt = pipeline.resolve(SESSIONS, "revamp", at="2")
    assert target["title"] == "Middle"
    assert prev["title"] == "Intro" and nxt["title"] == "Finale"


def test_check_revisable_refuses_a_revamp_with_nothing_to_revise():
    from tr import pipeline

    target = {"seq": "2", "title": "Middle", "tr_doc": None}
    label = lambda s: f"2. {s['title']}"

    try:
        pipeline.check_revisable("revamp", target, None, label)
    except SystemExit as e:
        assert "no TR doc to revise" in str(e)
        assert "mode 'new'" in str(e), "must say what to do instead"
    else:
        raise AssertionError("a revamp with no source doc must be refused")

    # a real doc passes, and so does mode new regardless
    pipeline.check_revisable("revamp", target, "# some doc", label)
    pipeline.check_revisable("new", None, None, label)


def _app_function(name):
    """Pull one function out of app.py without importing it.

    app.py issues Streamlit calls at module scope, so it cannot be imported in a
    plain test process.
    """
    src = (Path(__file__).resolve().parent.parent / "app.py").read_text(
        encoding="utf-8")
    start = src.index(f"def {name}")
    end = src.index("\n\n\n", start)
    ns = {}
    exec(src[start:end], ns)
    return ns[name]


def test_blocking_reason_always_returns_a_real_bool():
    """Regression: `target and not target.get(...)` returns None when target is
    None, and Streamlit's disabled= goes into a protobuf field that rejects it.
    A None here crashed the app on every 'new' mode run."""
    blocking_reason = _app_function("blocking_reason")

    cases = [
        ("", True, None, True),                        # no topic
        ("T", False, None, True),                      # no api key
        ("T", True, None, False),                      # new mode, ready
        ("T", True, {"tr_doc": "a.md"}, False),        # revamp with a doc
        ("T", True, {"tr_doc": None}, True),           # revamp, no doc
        ("T", True, {}, True),                         # revamp, key absent
    ]
    for topic, has_key, target, want in cases:
        blocked, why = blocking_reason(topic, has_key, target)
        assert blocked is want, f"{(topic, has_key, target)} -> {blocked!r}"
        assert type(blocked) is bool, \
            f"must be a real bool, got {type(blocked).__name__} for {target!r}"
        assert bool(why) == blocked, "a blocked button must say why"


def test_subtopics_reach_the_prompt_and_carry_the_scope_rules():
    from tr import pipeline

    spec = {"topic": "T", "slug": "t", "mode": "new", "position": "middle",
            "after": "2", "subtopics": "- A\n- B\n- C"}
    ctx = pipeline.build_context(spec, "Course", None, SESSIONS[0], SESSIONS[1],
                                 lambda s: s["title"])
    assert ctx["subtopics"] == "- A\n- B\n- C", "must survive into context.json"

    # The prompt is hard-wrapped, so any phrase can straddle a line break.
    # Assert against a whitespace-collapsed view, or a reflow breaks the tests
    # without anything actually being wrong.
    flat = " ".join(generate.prompt("tr_doc.md").split())
    for phrase in (
        "REQUIRED SUB-TOPICS",
        "minimum, not a ceiling",
        "ADDED BEYOND SCOPE",
        # the list is a coverage contract, not a running order
        "coverage contract, not a sequence",
        "The order is yours to choose",
        "listed order carries no meaning",
        "never introduce one that depends on a later one",
        "must lead into whichever sub-topic you place first",
        "Do not explain or justify your ordering",
    ):
        assert phrase in flat, f"scope rule missing: {phrase!r}"
    # and the no-subtopics path is stated, so an empty list is not ambiguous
    # the no-subtopics path is stated too, so an empty list is not ambiguous
    assert "If no sub-topics were supplied" in flat


def test_trailing_blocks_recognise_every_report_heading():
    """The writer's report blocks must not leak into the learner-facing doc,
    whether or not it prefixes them with markdown hashes."""
    for heading in ("SOURCE ISSUES", "OPEN MARKERS", "CHANGES MADE",
                    "ADDED BEYOND SCOPE"):
        for prefix in ("", "## ", "### "):
            doc, trailing = generate.split_trailing_blocks(
                f"# Doc\nreal content here\n\n{prefix}{heading}\n- something\n")
            assert doc == "# Doc\nreal content here", f"{prefix}{heading}"
            assert heading in trailing
            assert "something" in trailing

    # a doc with no report blocks is returned whole
    doc, trailing = generate.split_trailing_blocks("# Doc\njust content")
    assert doc == "# Doc\njust content" and trailing == ""

    # a heading mentioned mid-prose must not truncate the doc
    body = "# Doc\nWe list SOURCE ISSUES inline here as prose.\nMore content.\n"
    doc, trailing = generate.split_trailing_blocks(body)
    assert "More content." in doc


def test_trust_matching():
    trusted = ["docs.anthropic.com", "arxiv.org"]
    assert research.is_trusted("https://docs.anthropic.com/en/api", trusted)
    assert research.is_trusted("https://arxiv.org/abs/2301.001", trusted)
    assert not research.is_trusted("https://evil-docs.anthropic.com.attacker.io/x", trusted)
    assert not research.is_trusted("https://medium.com/@someone/post", trusted)
    assert not research.is_trusted("https://notarxiv.org/abs/1", trusted)


def test_trust_matching_rejects_typosquats_seen_in_the_wild():
    """A live research run returned developers-openai.com (hyphen, not dot)
    alongside the genuine developers.openai.com. A doc citing the squat as
    official would be worse than citing nothing."""
    trusted = corpus.trusted_domains("sources.yaml")

    assert research.is_trusted("https://developers.openai.com/api/docs", trusted)
    assert not research.is_trusted("https://developers-openai.com/docs/x", trusted)

    # the same trick on the other listed providers
    for squat in ("https://docs-anthropic.com/x",
                  "https://pydantic-dev.io/x",
                  "https://github-com.example.net/x",
                  "https://huggingface-co.net/x",
                  "https://openai.com.evil.test/x"):
        assert not research.is_trusted(squat, trusted), squat

    # and the real base domains still pass, including subdomains
    for good in ("https://openai.com/index/x",
                 "https://platform.openai.com/docs",
                 "https://docs.pydantic.dev/latest/",
                 "https://pydantic.dev/docs/x",
                 "https://pypi.org/project/pydantic/"):
        assert research.is_trusted(good, trusted), good


def test_downstream_flag_between():
    nxt = {"seq": "5", "unit_number": "45", "title": "Agents", "tr_doc": "c.md"}
    msg = generate.downstream_flag(nxt, "Structured Outputs")
    assert "Agents" in msg and "c.md" in msg and "Structured Outputs" in msg
    assert "5" in msg


def test_downstream_flag_last_has_nothing_downstream():
    msg = generate.downstream_flag(None, "X", position="last")
    assert "end of the course" in msg
    assert "c.md" not in msg


def test_downstream_flag_first_names_the_displaced_opener():
    nxt = {"seq": "1", "title": "Your Learning Journey", "tr_doc": "a.md"}
    msg = generate.downstream_flag(nxt, "Course Primer", position="first")
    assert "first session of the course" in msg
    assert "Your Learning Journey" in msg and "a.md" in msg


def test_downstream_flag_revamp_reports_no_insertion():
    target = {"seq": "12", "title": "Prompt Engineering", "tr_doc": "p.md"}
    nxt = {"seq": "13", "title": "Image Generation", "tr_doc": "n.md"}
    msg = generate.downstream_flag(nxt, "Prompt Engineering", mode="revamp",
                                   target=target)
    assert "No session was inserted" in msg
    assert "course order is unchanged" in msg
    assert "p.md" in msg, "must name the file this output would replace"
    assert "does not overwrite" in msg
    # the next session is only a soft check, never a mandatory rewrite
    assert "may no longer hold" in msg

    # with no next session there is nothing even to soft-check
    solo = generate.downstream_flag(None, "X", mode="repurpose", target=target)
    assert "No session was inserted" in solo
    assert "may no longer hold" not in solo


def test_fill_marks_missing_values():
    filled = generate.fill("Course: {{course}} / Level: {{level}}",
                           {"course": "Gen AI", "level": None})
    assert "Gen AI" in filled
    assert "[NEEDS: level]" in filled


def test_split_trailing_blocks():
    doc, trailing = generate.split_trailing_blocks(
        "# Doc\nbody\n\nSOURCE ISSUES\n- version mismatch\n")
    assert doc == "# Doc\nbody"
    assert "version mismatch" in trailing
    doc, trailing = generate.split_trailing_blocks("# Doc\nbody")
    assert trailing == ""


def test_slugify():
    assert slugify("Structured Outputs with Pydantic!") == "structured-outputs-with-pydantic"
    assert slugify("RAG  --  Part 2") == "rag-part-2"


if __name__ == "__main__":
    tests = [v for k, v in sorted(globals().items()) if k.startswith("test_")]
    failed = 0
    for t in tests:
        try:
            t()
            print(f"PASS  {t.__name__}")
        except Exception as e:
            failed += 1
            print(f"FAIL  {t.__name__}: {type(e).__name__}: {e}")
    print(f"\n{len(tests) - failed}/{len(tests)} passed")
    sys.exit(1 if failed else 0)
